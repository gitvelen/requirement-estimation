import os
from pathlib import Path

from docx import Document
from pptx import Presentation
from reportlab.pdfgen import canvas

BASE_DIR = Path(__file__).resolve().parent.parent / "samples" / "v4"
EVIDENCE_DIR = BASE_DIR / "evidence"
ESB_DIR = BASE_DIR / "esb"
CODE_DIR = BASE_DIR / "spring_boot_demo" / "src" / "main" / "java" / "com" / "example" / "demo"


def ensure_dirs():
    for path in [EVIDENCE_DIR, ESB_DIR, CODE_DIR]:
        path.mkdir(parents=True, exist_ok=True)


def write_esb_csv():
    content = (
        "提供方系统简称,提供方中文名称,调用方系统简称,调用方中文名称,交易名称,交易码,服务场景码,状态\n"
        "SYS_A,系统A,SYS_B,系统B,白名单查询,SOFP500001,SC001,正常使用\n"
        "SYS_A,系统A,SYS_C,系统C,额度同步,SOFP500002,SC002,正常使用\n"
        "SYS_D,系统D,SYS_A,系统A,授信校验,SOFP500003,SC003,废弃使用\n"
    )
    (ESB_DIR / "esb_sample.csv").write_text(content, encoding="utf-8")


def write_docx():
    doc = Document()
    doc.add_heading("系统A 证据材料", level=0)
    doc.add_paragraph("系统名称: 系统A (SYS_A)")
    doc.add_paragraph("核心功能: 白名单查询、额度同步")
    doc.add_paragraph("集成点: ESB交易 SOFP500001/500002")
    doc.add_paragraph("关键约束: 日终批量同步")
    doc.save(EVIDENCE_DIR / "evidence_sample.docx")


def write_pptx():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "系统A 架构概览"
    body = slide.placeholders[1].text_frame
    body.text = "系统A 主要模块"
    body.add_paragraph().text = "- 账户管理"
    body.add_paragraph().text = "- 额度管理"
    body.add_paragraph().text = "- 对接ESB"
    prs.save(EVIDENCE_DIR / "evidence_sample.pptx")


def write_pdf():
    path = EVIDENCE_DIR / "evidence_sample.pdf"
    c = canvas.Canvas(str(path))
    c.setFont("Helvetica", 12)
    c.drawString(50, 800, "系统A 证据材料 (PDF)")
    c.drawString(50, 780, "系统名称: 系统A (SYS_A)")
    c.drawString(50, 760, "功能: 白名单查询、额度同步")
    c.drawString(50, 740, "集成点: ESB交易 SOFP500001/500002")
    c.save()


def write_java_sources():
    controller = """
package com.example.demo;

import org.springframework.web.bind.annotation.*;

@RestController
@RequestMapping("/api/demo")
public class DemoController {
    @GetMapping("/hello")
    public String hello() {
        return "hello";
    }

    @PostMapping("/submit")
    public String submit(@RequestBody String payload) {
        return payload;
    }
}
""".strip()

    scheduled = """
package com.example.demo;

import org.springframework.scheduling.annotation.Scheduled;
import org.springframework.stereotype.Component;

@Component
public class SyncJob {
    @Scheduled(cron = "0 0/5 * * * ?")
    public void syncQuota() {
        // sync logic
    }
}
""".strip()

    listener = """
package com.example.demo;

import org.springframework.kafka.annotation.KafkaListener;
import org.springframework.stereotype.Component;

@Component
public class EventListener {
    @KafkaListener(topics = "demo-topic")
    public void handle(String message) {
        // handle event
    }
}
""".strip()

    outbound = """
package com.example.demo;

import org.springframework.cloud.openfeign.FeignClient;
import org.springframework.web.bind.annotation.GetMapping;

@FeignClient(name = "user-service")
public interface UserClient {
    @GetMapping("/api/user/profile")
    String getProfile();
}
""".strip()

    (CODE_DIR / "DemoController.java").write_text(controller, encoding="utf-8")
    (CODE_DIR / "SyncJob.java").write_text(scheduled, encoding="utf-8")
    (CODE_DIR / "EventListener.java").write_text(listener, encoding="utf-8")
    (CODE_DIR / "UserClient.java").write_text(outbound, encoding="utf-8")


def write_readme():
    content = """
# V4样例数据

- evidence/: 证据材料样例（DOCX/PPTX/PDF）
- esb/esb_sample.csv: ESB明细样例
- spring_boot_demo/: Spring Boot扫描样例（含Controller/Scheduled/Listener/Feign）
""".strip() + "\n"
    (BASE_DIR / "README.md").write_text(content, encoding="utf-8")


def main():
    ensure_dirs()
    write_esb_csv()
    write_docx()
    write_pptx()
    write_pdf()
    write_java_sources()
    write_readme()


if __name__ == "__main__":
    main()
